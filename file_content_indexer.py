#!/usr/bin/env python3
"""
File Content Indexer for Discord Indexer

This module downloads files from Discord and extracts their text content
for full-text search capabilities. Supports various file formats including:
- PDF files
- Microsoft Office documents (Word, Excel, PowerPoint)
- Text files
- Code files
- And more
"""

import os
import sqlite3
import requests
import tempfile
import logging
import mimetypes
from datetime import datetime
from pathlib import Path
import hashlib
import json
import io

# Text extraction libraries
try:
    import PyPDF2
    import pdfplumber
except ImportError:
    PyPDF2 = None
    pdfplumber = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import pytesseract
    from PIL import Image
except ImportError:
    pytesseract = None
    Image = None

try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

class FileContentIndexer:
    def __init__(self, db_path, download_dir="downloads"):
        self.db_path = db_path
        self.download_dir = Path(download_dir)
        self.download_dir.mkdir(exist_ok=True)
        
        # Setup logging
        self.logger = logging.getLogger('FileContentIndexer')
        
        # Supported file types for content extraction
        self.supported_types = {
            'application/pdf': self._extract_pdf_content,
            'application/msword': self._extract_docx_content,
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document': self._extract_docx_content,
            'application/vnd.ms-excel': self._extract_excel_content,
            'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': self._extract_excel_content,
            'application/vnd.ms-powerpoint': self._extract_pptx_content,
            'application/vnd.openxmlformats-officedocument.presentationml.presentation': self._extract_pptx_content,
            'text/plain': self._extract_text_content,
            'text/markdown': self._extract_text_content,
            'text/csv': self._extract_text_content,
            'application/json': self._extract_text_content,
            'application/xml': self._extract_text_content,
            'text/xml': self._extract_text_content,
            'text/html': self._extract_text_content,
            'text/css': self._extract_text_content,
            'text/javascript': self._extract_text_content,
            'application/javascript': self._extract_text_content,
            'text/x-python': self._extract_text_content,
            'text/x-java-source': self._extract_text_content,
            'text/x-c': self._extract_text_content,
            'text/x-c++': self._extract_text_content,
            # Image types for OCR
            'image/jpeg': self._extract_image_content,
            'image/jpg': self._extract_image_content,
            'image/png': self._extract_image_content,
            'image/gif': self._extract_image_content,
            'image/bmp': self._extract_image_content,
            'image/tiff': self._extract_image_content,
            'image/webp': self._extract_image_content,
        }
        
        # Initialize database
        self._init_database()
    
    def _init_database(self):
        """Initialize the file_content table if it doesn't exist"""
        conn = sqlite3.connect(self.db_path)
        cursor = conn.cursor()
        
        cursor.execute('''
            CREATE TABLE IF NOT EXISTS file_content (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                file_id INTEGER NOT NULL,
                content_text TEXT,
                extraction_method TEXT,
                extracted_at TEXT NOT NULL,
                file_size_bytes INTEGER,
                FOREIGN KEY (file_id) REFERENCES indexed_files (id) ON DELETE CASCADE,
                UNIQUE(file_id)
            )
        ''')
        
        # Create index for faster searches
        cursor.execute('''
            CREATE INDEX IF NOT EXISTS idx_file_content_text 
            ON file_content(content_text)
        ''')
        
        # Index removed - content_hash column no longer exists
        
        conn.commit()
        conn.close()
    
    def get_connection(self):
        """Get database connection"""
        conn = sqlite3.connect(self.db_path, timeout=30.0)
        conn.execute('PRAGMA journal_mode=WAL')
        conn.execute('PRAGMA busy_timeout=30000')
        return conn
    
    def download_file(self, file_url, filename):
        """Download file from URL and return local path"""
        try:
            response = requests.get(file_url, timeout=30)
            response.raise_for_status()
            
            # Create safe filename
            safe_filename = "".join(c for c in filename if c.isalnum() or c in (' ', '.', '_', '-')).rstrip()
            if not safe_filename:
                safe_filename = f"file_{hashlib.md5(file_url.encode()).hexdigest()[:8]}"
            
            file_path = self.download_dir / safe_filename
            
            # If file exists, add timestamp to make it unique
            if file_path.exists():
                stem = file_path.stem
                suffix = file_path.suffix
                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                file_path = self.download_dir / f"{stem}_{timestamp}{suffix}"
            
            with open(file_path, 'wb') as f:
                f.write(response.content)
            
            return file_path
            
        except Exception as e:
            self.logger.error(f"Error downloading file {file_url}: {str(e)}")
            return None
    
    def _extract_pdf_content(self, file_path):
        """Extract text content from PDF files with OCR for images"""
        content = ""
        extraction_method = "none"
        
        # Try pdfplumber first (better for complex layouts)
        if pdfplumber:
            try:
                with pdfplumber.open(file_path) as pdf:
                    for page in pdf.pages[:50]:  # Limit to first 50 pages
                        page_text = page.extract_text()
                        if page_text:
                            content += page_text + "\n"
                if content.strip():
                    extraction_method = "pdfplumber"
            except Exception as e:
                self.logger.warning(f"pdfplumber failed for {file_path}: {str(e)}")
        
        # Fallback to PyPDF2 if pdfplumber didn't work
        if not content.strip() and PyPDF2:
            try:
                with open(file_path, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    for page_num in range(min(50, len(pdf_reader.pages))):
                        page = pdf_reader.pages[page_num]
                        content += page.extract_text() + "\n"
                if content.strip():
                    extraction_method = "PyPDF2"
            except Exception as e:
                self.logger.warning(f"PyPDF2 failed for {file_path}: {str(e)}")
        
        # If no text extracted or very little text, try OCR on images in PDF
        if (not content.strip() or len(content.strip()) < 100) and fitz and pytesseract and Image:
            try:
                ocr_content = self._extract_pdf_images_ocr(file_path)
                if ocr_content:
                    content += "\n" + ocr_content
                    extraction_method = f"{extraction_method}+ocr" if extraction_method != "none" else "pdf-ocr"
            except Exception as e:
                self.logger.warning(f"PDF OCR failed for {file_path}: {str(e)}")
        
        return content, extraction_method if extraction_method != "none" else "none"
    
    def _extract_pdf_images_ocr(self, file_path):
        """Extract and OCR images from PDF using PyMuPDF"""
        ocr_content = ""
        
        try:
            doc = fitz.open(file_path)
            
            for page_num in range(min(50, len(doc))):  # Limit to first 50 pages
                page = doc[page_num]
                image_list = page.get_images()
                
                for img_index, img in enumerate(image_list):
                    try:
                        # Get the image
                        xref = img[0]
                        pix = fitz.Pixmap(doc, xref)
                        
                        # Convert to PIL Image if not GRAY or RGB
                        if pix.n - pix.alpha < 4:  # GRAY or RGB
                            img_data = pix.tobytes("ppm")
                            pil_image = Image.open(io.BytesIO(img_data))
                        else:  # CMYK: convert to RGB first
                            pix1 = fitz.Pixmap(fitz.csRGB, pix)
                            img_data = pix1.tobytes("ppm")
                            pil_image = Image.open(io.BytesIO(img_data))
                            pix1 = None
                        
                        pix = None
                        
                        # Skip very small images (likely decorative)
                        if pil_image.width < 50 or pil_image.height < 50:
                            continue
                        
                        # Perform OCR
                        text = pytesseract.image_to_string(pil_image, lang='eng')
                        if text.strip():
                            ocr_content += f"\n[Image {page_num+1}-{img_index+1}]: {text.strip()}\n"
                        
                    except Exception as e:
                        self.logger.warning(f"Error processing image {img_index} on page {page_num}: {str(e)}")
                        continue
            
            doc.close()
            
        except Exception as e:
            self.logger.error(f"Error extracting images from PDF {file_path}: {str(e)}")
        
        return ocr_content
    
    def _extract_docx_content(self, file_path):
        """Extract text content from Word documents with OCR for images"""
        if not Document:
            return "", "none"
        
        try:
            doc = Document(file_path)
            content = ""
            extraction_method = "python-docx"
            
            # Extract paragraphs
            for paragraph in doc.paragraphs:
                content += paragraph.text + "\n"
            
            # Extract tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        content += cell.text + " "
                    content += "\n"
            
            # Extract OCR content from embedded images
            if pytesseract and Image:
                try:
                    ocr_content = self._extract_docx_images_ocr(file_path)
                    if ocr_content:
                        content += "\n" + ocr_content
                        extraction_method = "python-docx+ocr"
                except Exception as e:
                    self.logger.warning(f"DOCX OCR failed for {file_path}: {str(e)}")
            
            return content, extraction_method
            
        except Exception as e:
            self.logger.error(f"Error extracting DOCX content from {file_path}: {str(e)}")
            return "", "error"
    
    def _extract_docx_images_ocr(self, file_path):
        """Extract and OCR images from Word documents"""
        ocr_content = ""
        
        try:
            import zipfile
            
            # Word documents are ZIP files
            with zipfile.ZipFile(file_path, 'r') as docx_zip:
                # Look for images in the media folder
                image_files = [f for f in docx_zip.namelist() if f.startswith('word/media/') and 
                             f.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff'))]
                
                for img_index, img_file in enumerate(image_files):
                    try:
                        # Extract image data
                        img_data = docx_zip.read(img_file)
                        pil_image = Image.open(io.BytesIO(img_data))
                        
                        # Skip very small images (likely decorative)
                        if pil_image.width < 50 or pil_image.height < 50:
                            continue
                        
                        # Perform OCR
                        text = pytesseract.image_to_string(pil_image, lang='eng')
                        if text.strip():
                            ocr_content += f"\n[Image {img_index+1}]: {text.strip()}\n"
                        
                    except Exception as e:
                        self.logger.warning(f"Error processing image {img_file}: {str(e)}")
                        continue
            
        except Exception as e:
            self.logger.error(f"Error extracting images from DOCX {file_path}: {str(e)}")
        
        return ocr_content
    
    def _extract_excel_content(self, file_path):
        """Extract text content from Excel files"""
        if not openpyxl:
            return "", "none"
        
        try:
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            content = ""
            
            for sheet_name in workbook.sheetnames[:10]:  # Limit to first 10 sheets
                sheet = workbook[sheet_name]
                content += f"Sheet: {sheet_name}\n"
                
                for row in sheet.iter_rows(max_row=1000, values_only=True):  # Limit to 1000 rows
                    row_text = [str(cell) if cell is not None else "" for cell in row]
                    content += " ".join(row_text) + "\n"
            
            return content, "openpyxl"
            
        except Exception as e:
            self.logger.error(f"Error extracting Excel content from {file_path}: {str(e)}")
            return "", "error"
    
    def _extract_pptx_content(self, file_path):
        """Extract text content from PowerPoint files with OCR for images"""
        if not Presentation:
            return "", "none"
        
        try:
            prs = Presentation(file_path)
            content = ""
            extraction_method = "python-pptx"
            
            for slide_num, slide in enumerate(prs.slides[:100]):  # Limit to 100 slides
                content += f"Slide {slide_num + 1}:\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        content += shape.text + "\n"
                
                content += "\n"
            
            # Extract OCR content from embedded images
            if pytesseract and Image:
                try:
                    ocr_content = self._extract_pptx_images_ocr(file_path)
                    if ocr_content:
                        content += "\n" + ocr_content
                        extraction_method = "python-pptx+ocr"
                except Exception as e:
                    self.logger.warning(f"PPTX OCR failed for {file_path}: {str(e)}")
            
            return content, extraction_method
            
        except Exception as e:
            self.logger.error(f"Error extracting PPTX content from {file_path}: {str(e)}")
            return "", "error"
    
    def _extract_pptx_images_ocr(self, file_path):
        """Extract and OCR images from PowerPoint files"""
        ocr_content = ""
        
        try:
            import zipfile
            
            # PowerPoint documents are ZIP files
            with zipfile.ZipFile(file_path, 'r') as pptx_zip:
                # Look for images in the media folder
                image_files = [f for f in pptx_zip.namelist() if f.startswith('ppt/media/') and 
                             f.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff'))]
                
                for img_index, img_file in enumerate(image_files):
                    try:
                        # Extract image data
                        img_data = pptx_zip.read(img_file)
                        pil_image = Image.open(io.BytesIO(img_data))
                        
                        # Skip very small images (likely decorative)
                        if pil_image.width < 50 or pil_image.height < 50:
                            continue
                        
                        # Perform OCR
                        text = pytesseract.image_to_string(pil_image, lang='eng')
                        if text.strip():
                            ocr_content += f"\n[Image {img_index+1}]: {text.strip()}\n"
                        
                    except Exception as e:
                        self.logger.warning(f"Error processing image {img_file}: {str(e)}")
                        continue
            
        except Exception as e:
            self.logger.error(f"Error extracting images from PPTX {file_path}: {str(e)}")
        
        return ocr_content
    
    def _extract_text_content(self, file_path):
        """Extract content from text-based files"""
        try:
            # Try different encodings
            encodings = ['utf-8', 'utf-16', 'latin-1', 'cp1252']
            
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        content = f.read(100000)  # Limit to 100KB
                    return content, f"text-{encoding}"
                except UnicodeDecodeError:
                    continue
            
            # If all encodings fail, try binary mode and decode with errors='ignore'
            with open(file_path, 'rb') as f:
                raw_content = f.read(100000)
                content = raw_content.decode('utf-8', errors='ignore')
            
            return content, "text-binary"
            
        except Exception as e:
            self.logger.error(f"Error extracting text content from {file_path}: {str(e)}")
            return "", "error"
    
    def _extract_image_content(self, file_path):
        """Extract text content from images using OCR"""
        if not pytesseract or not Image:
            return "", "none"
        
        try:
            # Open and process the image
            with Image.open(file_path) as img:
                # Convert to RGB if necessary (for better OCR results)
                if img.mode != 'RGB':
                    img = img.convert('RGB')
                
                # Extract text using Tesseract OCR
                # Use custom config for better accuracy
                custom_config = r'--oem 3 --psm 6 -c tessedit_char_whitelist=0123456789ABCDEFGHIJKLMNOPQRSTUVWXYZabcdefghijklmnopqrstuvwxyz !"#$%&\'()*+,-./:;<=>?@[\]^_`{|}~'
                
                try:
                    # Try with custom config first
                    content = pytesseract.image_to_string(img, config=custom_config)
                except Exception:
                    # Fallback to default config
                    content = pytesseract.image_to_string(img)
                
                # Clean up the extracted text
                content = content.strip()
                
                if content:
                    self.logger.info(f"OCR extracted {len(content)} characters from image {file_path}")
                    return content, "tesseract-ocr"
                else:
                    self.logger.info(f"No text found in image {file_path}")
                    return "", "no-text"
                    
        except Exception as e:
            self.logger.error(f"Error extracting text from image {file_path}: {str(e)}")
            return "", "error"
    
    def extract_file_content(self, file_id, file_url, filename, file_type):
        """Extract content from a file and store it in the database"""
        try:
            # Check if content already exists
            conn = self.get_connection()
            cursor = conn.cursor()
            
            cursor.execute("SELECT id FROM file_content WHERE file_id = ?", (file_id,))
            if cursor.fetchone():
                self.logger.info(f"Content already exists for file {file_id}")
                conn.close()
                return True
            
            conn.close()
            
            # Check if file type is supported
            if file_type not in self.supported_types:
                self.logger.info(f"File type {file_type} not supported for content extraction")
                return False
            
            # Download file
            self.logger.info(f"Downloading file {filename} for content extraction")
            file_path = self.download_file(file_url, filename)
            
            if not file_path or not file_path.exists():
                self.logger.error(f"Failed to download file {filename}")
                return False
            
            try:
                # Extract content
                extract_func = self.supported_types[file_type]
                content, method = extract_func(file_path)
                
                if not content.strip():
                    self.logger.warning(f"No content extracted from {filename}")
                    return False
                
                # Content extracted successfully
                
                # Get file size
                file_size = file_path.stat().st_size
                
                # Store in database
                conn = self.get_connection()
                cursor = conn.cursor()
                
                cursor.execute('''
                    INSERT OR REPLACE INTO file_content 
                    (file_id, content_text, extraction_method, extracted_at, file_size_bytes)
                    VALUES (?, ?, ?, ?, ?)
                ''', (file_id, content, method, datetime.now().isoformat(), file_size))
                
                conn.commit()
                conn.close()
                
                self.logger.info(f"Successfully extracted and stored content for {filename} ({len(content)} chars)")
                return True
                
            finally:
                # Clean up downloaded file
                try:
                    file_path.unlink()
                except Exception as e:
                    self.logger.warning(f"Failed to delete temporary file {file_path}: {str(e)}")
            
        except Exception as e:
            self.logger.error(f"Error extracting content from file {filename}: {str(e)}")
            return False
    
    def index_all_files(self):
        """Index content for all files that don't have content yet"""
        conn = self.get_connection()
        cursor = conn.cursor()
        
        # Get files without content
        cursor.execute('''
            SELECT f.id, f.filename, f.file_url, f.file_type
            FROM indexed_files f
            LEFT JOIN file_content fc ON f.id = fc.file_id
            WHERE fc.file_id IS NULL
            ORDER BY f.indexed_at DESC
        ''')
        
        files_to_process = cursor.fetchall()
        conn.close()
        
        self.logger.info(f"Found {len(files_to_process)} files to process for content extraction")
        
        success_count = 0
        for file_id, filename, file_url, file_type in files_to_process:
            if self.extract_file_content(file_id, file_url, filename, file_type):
                success_count += 1
        
        self.logger.info(f"Successfully processed {success_count}/{len(files_to_process)} files")
        return success_count, len(files_to_process)
    
    def search_file_content(self, query, limit=50):
        """Search through file content"""
        conn = self.get_connection()
        cursor = conn.cursor()
        
        # Simple text search (can be enhanced with FTS later)
        cursor.execute('''
            SELECT f.id, f.filename, f.file_url, f.file_type, f.author_name, 
                   f.channel_name, f.timestamp, fc.content_text
            FROM indexed_files f
            JOIN file_content fc ON f.id = fc.file_id
            WHERE fc.content_text LIKE ?
            ORDER BY f.timestamp DESC
            LIMIT ?
        ''', (f'%{query}%', limit))
        
        results = cursor.fetchall()
        conn.close()
        
        return results
    
    def get_content_stats(self):
        """Get statistics about indexed content"""
        conn = self.get_connection()
        cursor = conn.cursor()
        
        # Total files with content
        cursor.execute("SELECT COUNT(*) FROM file_content")
        files_with_content = cursor.fetchone()[0]
        
        # Total files
        cursor.execute("SELECT COUNT(*) FROM indexed_files")
        total_files = cursor.fetchone()[0]
        
        # Content by extraction method
        cursor.execute('''
            SELECT extraction_method, COUNT(*) 
            FROM file_content 
            GROUP BY extraction_method
        ''')
        methods = dict(cursor.fetchall())
        
        # Total content size
        cursor.execute("SELECT SUM(LENGTH(content_text)) FROM file_content")
        total_content_size = cursor.fetchone()[0] or 0
        
        conn.close()
        
        return {
            'files_with_content': files_with_content,
            'total_files': total_files,
            'coverage_percentage': (files_with_content / total_files * 100) if total_files > 0 else 0,
            'extraction_methods': methods,
            'total_content_characters': total_content_size
        }

if __name__ == '__main__':
    # Example usage
    indexer = FileContentIndexer('indexer.db')
    success, total = indexer.index_all_files()
    print(f"Processed {success}/{total} files")
    
    stats = indexer.get_content_stats()
    print(f"Content statistics: {json.dumps(stats, indent=2)}")